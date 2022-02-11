from django.shortcuts import render

from pptx import Presentation
from pptx.util import Inches

from .models import Employee

# Create your views here.

from django.http import HttpResponse
from django.shortcuts import get_object_or_404, render

from .models import Employee

from thumbnail import generate_thumbnail

def index(request):
	employee_list = Employee.objects.all()
	return render(request, 'emplist/list.html', {'employee_list': employee_list })

def detail(request, employee_id):

	emp = get_object_or_404(Employee, pk=employee_id)
	return render(request, 'emplist/detail.html', {'employee': emp })

def genpptx(reqeust, employee_id):

	emp = Employee.objects.get( pk=employee_id )

	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	shapes = slide.shapes

	shapes.title.text = 'Sample Table'

	rows = 4
	cols = 2
	left = top = Inches(2.0)
	width = Inches(6.0)
	height = Inches(0.8)

	table = shapes.add_table(rows, cols, left, top, width, height).table

	table.columns[0].width = Inches(2.0)
	table.columns[1].width = Inches(4.0)

	table.cell(0,0).text = 'First Name'
	table.cell(0,1).text = emp.firstname_text

	table.cell(1,0).text = 'Last Name'
	table.cell(1,1).text = emp.lastname_text

	table.cell(2,0).text = 'Address'
	table.cell(2,1).text = emp.address_text

	table.cell(3,0).text = 'Phone'
	table.cell(3,1).text = emp.phone_text

	prs.save('test.pptx')

	######################
	# Generate a thumbnail
	######################

	options = {
		'trim': False,
		'height': 300,
		'width': 300,
		'quality': 85,
		'type': 'thumbnail'
	}

	generate_thumbnail('test.pptx', 'thumbnail.png', options )

	file = open("test.pptx", "rb")

	response = HttpResponse(file.read(), content_type="application/pptx")
	response['Content-Disposition'] = 'attachment; filename=%s' % emp.firstname_text + '.pptx'
	return response

def regist(request):

	firstname = request.POST['firstname']
	lastname = request.POST['lastname']
	address = request.POST['address']
	phone = request.POST['phone']

	emp = Employee( firstname_text=firstname, lastname_text=lastname, address_text=address, phone_text = phone )
	emp.save()

	employee_list = Employee.objects.all()
	return render(request, 'emplist/list.html', {'employee_list': employee_list })



		
def registform(request):

	return render(request, 'emplist/registform.html')

